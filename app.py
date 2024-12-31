import os
import pyttsx3
from flask import Flask, request, jsonify, render_template
from werkzeug.utils import secure_filename
import pdfplumber
from pptx import Presentation
import requests
from bs4 import BeautifulSoup

app = Flask(__name__)

UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'pdf', 'pptx'}
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def speak_text(text, callback):
    engine = pyttsx3.init()
    engine.setProperty('rate', 150)  # Adjust speed
    engine.setProperty('volume', 1)  # Adjust volume (0.0 to 1.0)

    def on_word(name, location, length):
        word = text[location:location + length]
        callback(word)  # Send the word to the frontend for highlighting

    engine.connect('started-word', on_word)
    engine.say(text)
    engine.runAndWait()

def extract_pdf_content(file_path):
    content = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            content += page.extract_text()  # Extract text from each page
    return content

def extract_ppt_content(file_path):
    content = ""
    presentation = Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                content += shape.text
    return content

def extract_website_content(url):
    response = requests.get(url)
    soup = BeautifulSoup(response.text, 'html.parser')
    content = " ".join([p.get_text() for p in soup.find_all('p')])  # Extract text from <p> tags
    return content

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/extract', methods=['POST'])
def extract():
    file = request.files.get('file')
    url = request.form.get('url')

    if file:
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(file_path)

            if filename.lower().endswith('.pdf'):
                content = extract_pdf_content(file_path)
            elif filename.lower().endswith('.pptx'):
                content = extract_ppt_content(file_path)

            return jsonify({"content": content})
    
    elif url:
        content = extract_website_content(url)
        return jsonify({"content": content})

    return jsonify({"message": "No valid file or URL provided"}), 400

@app.route('/talk', methods=['POST'])
def talk():
    content = request.form.get('content')
    if content:
        def callback(word):
            # Send the current word to be highlighted
            socketio.emit('highlight_word', {'word': word})
        
        speak_text(content, callback)  # Speak content and highlight text
        return jsonify({"message": "Content is being spoken!"})
    else:
        return jsonify({"message": "No content provided to speak."}), 400

if __name__ == '__main__':
    from flask_socketio import SocketIO
    socketio = SocketIO(app)
    socketio.run(app, debug=True)
